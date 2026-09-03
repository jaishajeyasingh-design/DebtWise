import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck(filename="FinShield_Pitch_Deck.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # -------------------------------------------------------------
    # Color Palette - Modern Fintech Theme
    # -------------------------------------------------------------
    NAVY_DARK    = RGBColor(10, 17, 40)       # #0A1128 Primary Dark Background
    CARD_BG      = RGBColor(18, 28, 58)       # #121C3A Card Background
    CARD_BORDER  = RGBColor(38, 54, 98)       # #263662 Card Border
    ACCENT_BLUE  = RGBColor(58, 134, 255)     # #3A86FF Electric Blue
    ACCENT_CYAN  = RGBColor(0, 240, 255)      # #00F0FF Vibrant Cyan
    ACCENT_TEAL  = RGBColor(16, 185, 129)     # #10B981 Emerald Green
    ACCENT_CORAL = RGBColor(244, 63, 94)      # #F43F5E Coral / Alert
    ACCENT_AMBER = RGBColor(245, 158, 11)     # #F59E0B Amber / Warning
    TEXT_WHITE   = RGBColor(255, 255, 255)
    TEXT_MUTED   = RGBColor(156, 163, 175)    # Slate 400
    TEXT_DIM     = RGBColor(100, 116, 139)    # Slate 500
    LIGHT_BG     = RGBColor(248, 250, 252)    # Light mode / white card
    
    # Blank slide layout
    blank_layout = prs.slide_layouts[6]
    
    # Helper: Add Solid Background
    def set_slide_background(slide, color=NAVY_DARK):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    # Helper: Add Header
    def add_header(slide, category, title, subtitle=None):
        # Category Pill / Supertitle
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11.7), Inches(0.35))
        tf = cat_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = category.upper()
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = ACCENT_CYAN
        p.font.name = "Segoe UI"
        
        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.78), Inches(11.7), Inches(0.6))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        p.font.name = "Segoe UI"
        
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.38), Inches(11.7), Inches(0.35))
            tf = sub_box.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(11)
            p.font.color.rgb = TEXT_MUTED
            p.font.name = "Segoe UI"

    # Helper: Create Card Shape
    def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        if border_color:
            card.line.color.rgb = border_color
            card.line.width = Pt(1.5)
        else:
            card.line.fill.background()
        return card

    # Helper: Add Speaker Notes
    def add_speaker_notes(slide, notes_text):
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = notes_text

    # =========================================================================
    # SLIDE 1: TITLE SLIDE
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide1, NAVY_DARK)
    
    # Decorative accent card
    accent_bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(0.15), Inches(3.8))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT_CYAN
    accent_bar.line.fill.background()

    # Title & Subtitle Box
    t_box = slide1.shapes.add_textbox(Inches(1.2), Inches(1.7), Inches(11.0), Inches(4.2))
    tf = t_box.text_frame
    tf.word_wrap = True
    
    # Badge
    p0 = tf.paragraphs[0]
    p0.text = "HACKATHON EDITION  |  RESPONSIBLE AI IN FINTECH"
    p0.font.size = Pt(11)
    p0.font.bold = True
    p0.font.color.rgb = ACCENT_CYAN
    p0.font.name = "Segoe UI"
    p0.space_after = Pt(14)
    
    # Main Product Name
    p1 = tf.add_paragraph()
    p1.text = "FinShield"
    p1.font.size = Pt(54)
    p1.font.bold = True
    p1.font.color.rgb = TEXT_WHITE
    p1.font.name = "Segoe UI"
    p1.space_after = Pt(8)
    
    # Value Proposition
    p2 = tf.add_paragraph()
    p2.text = "AI-Assisted Financial Distress Intervention Engine"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT_BLUE
    p2.font.name = "Segoe UI"
    p2.space_after = Pt(14)
    
    p3 = tf.add_paragraph()
    p3.text = "Moving banks from merely detecting high-risk customers to safely and sustainably resolving distress — before default occurs."
    p3.font.size = Pt(14)
    p3.font.color.rgb = TEXT_MUTED
    p3.font.name = "Segoe UI"
    p3.space_after = Pt(24)

    # 3 Key Feature Pills at bottom
    pills = [
        ("DIAGNOSE", "XGBoost + SHAP Root Cause Analysis"),
        ("CALCULATE", "Deterministic Repayment Capacity"),
        ("INTERVENE", "Safety-Filtered 7-Tier Action Ladder")
    ]
    for i, (tag, desc) in enumerate(pills):
        card = add_card(slide1, Inches(1.2 + i * 3.8), Inches(5.8), Inches(3.5), Inches(0.9), CARD_BG, CARD_BORDER)
        box = slide1.shapes.add_textbox(Inches(1.35 + i * 3.8), Inches(5.85), Inches(3.2), Inches(0.8))
        tf_pill = box.text_frame
        tf_pill.word_wrap = True
        p_t = tf_pill.paragraphs[0]
        p_t.text = tag
        p_t.font.size = Pt(10)
        p_t.font.bold = True
        p_t.font.color.rgb = ACCENT_CYAN
        p_d = tf_pill.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = TEXT_WHITE

    add_speaker_notes(slide1, """Good morning judges. We are thrilled to present FinShield.
Today, every major bank has sophisticated models to predict who is at risk of defaulting. But here is the critical gap: detection is not a solution. What happens AFTER a customer is flagged?
FinShield is an AI-assisted financial distress intervention engine that diagnoses the root cause of distress, computes sustainable affordability, safety-filters interventions, and resolves distress before it leads to default or financial exclusion.""")

    # =========================================================================
    # SLIDE 2: THE CORE PROBLEM — THE POST-FLAGGING INTERVENTION VOID
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide2)
    add_header(slide2, "The Industry Dilemma", "Banks Know Who Is High-Risk. The Real Crisis Is What Happens Next.", "Risk detection is mature, but the post-flagging intervention workflow remains broken, punitive, and manual.")

    # Left Card: Traditional Broken Cycle
    add_card(slide2, Inches(0.8), Inches(1.85), Inches(5.6), Inches(4.9), CARD_BG, ACCENT_CORAL)
    t_left = slide2.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.5))
    tf_l = t_left.text_frame
    tf_l.word_wrap = True
    
    p = tf_l.paragraphs[0]
    p.text = "THE CURRENT REALITY: THE DEFAULT ESCALATOR"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CORAL
    p.space_after = Pt(12)
    
    points_left = [
        ("Passive Risk Flagging", "Models output a probability score (PD). Nothing happens until missed payments accumulate."),
        ("Blind Collections Handoff", "Treated as a recovery problem rather than an affordability mismatch. Aggressive notices begin."),
        ("One-Size-Fits-All Penalties", "Late fees and overdraft charges are added, directly worsening the customer's cash-flow crunch."),
        ("Binary Outcome", "Leads directly to loan default, expensive legal recovery, write-offs, and severe financial exclusion.")
    ]
    for title, body in points_left:
        p_t = tf_l.add_paragraph()
        p_t.text = "•  " + title
        p_t.font.size = Pt(11)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_WHITE
        p_b = tf_l.add_paragraph()
        p_b.text = "    " + body
        p_b.font.size = Pt(10)
        p_b.font.color.rgb = TEXT_MUTED
        p_b.space_after = Pt(8)

    # Right Card: The FinShield Core Question
    add_card(slide2, Inches(6.8), Inches(1.85), Inches(5.7), Inches(4.9), CARD_BG, ACCENT_CYAN)
    t_right = slide2.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(5.3), Inches(4.5))
    tf_r = t_right.text_frame
    tf_r.word_wrap = True
    
    p = tf_r.paragraphs[0]
    p.text = "THE UNANSWERED INTERVENTION QUESTION"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.space_after = Pt(14)
    
    p_q = tf_r.add_paragraph()
    p_q.text = "“How can a bank responsibly understand WHY a customer is in distress, determine what they can sustainably afford, choose the safest intervention, and act early enough to prevent default?”"
    p_q.font.size = Pt(13)
    p_q.font.bold = True
    p_q.font.color.rgb = TEXT_WHITE
    p_q.space_after = Pt(16)
    
    pillars = [
        ("Root-Cause Clarity", "Distinguish income drops from timing mismatches or unexpected expense spikes."),
        ("Sustainable Capacity Floor", "Transparently calculate true affordability before offering any plan."),
        ("Safety & Friction Filter", "Reject predatory refinancing or debt traps that escalate customer burden."),
        ("Empathy & Customer Agency", "Provide choice, explainability, and consent instead of forced actions.")
    ]
    for title, body in pillars:
        p_t = tf_r.add_paragraph()
        p_t.text = "✔  " + title + ": "
        p_t.font.size = Pt(10)
        p_t.font.bold = True
        p_t.font.color.rgb = ACCENT_TEAL
        
        # Append inline
        p_t.text += body
        p_t.font.color.rgb = TEXT_MUTED
        p_t.space_after = Pt(6)

    add_speaker_notes(slide2, """When a bank's risk model flags a customer as high-risk, what happens next? Today, the answer is usually: wait for 60-90 days of missed payments, charge penalty fees, and send the file to aggressive collections.
This makes no sense. The penalty fees worsen the exact cash-flow issue causing the distress!
The missing piece is an intelligent, responsible intervention layer. FinShield steps in right at the moment of distress to ask: WHY are they struggling, WHAT can they sustainably afford, and WHAT is the safest treatment?""")

    # =========================================================================
    # SLIDE 3: KEY DIFFERENTIATOR — DIAGNOSIS BEFORE TREATMENT
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide3)
    add_header(slide3, "Strategic Positioning", "Diagnosis Before Treatment: The FinShield Paradigm", "We are not another risk scorer or chatbot. We are a clinical intervention engine for financial health.")

    # 4 Contrast Cards
    cards_data = [
        ("WHAT WE ARE NOT", [
            ("Credit-Risk Scorecard", "Scorecards tell you probability of default (PD), but never tell you how to cure it."),
            ("Banking Chatbot", "Conversational bots lack deterministic financial guardrails and risk engines."),
            ("Loan Recommender", "Pushing more credit to a distressed borrower creates predatory debt spirals."),
            ("Generic BI Dashboard", "Passive charts don't generate safety-checked, actionable intervention paths.")
        ], ACCENT_CORAL),
        ("WHAT FINSHIELD IS", [
            ("Root-Cause Diagnostic", "Identifies the exact financial mechanism driving distress (5 distinct archetypes)."),
            ("Deterministic Capacity Engine", "Calculates exact sustainable repayment capacity based on living cost floors."),
            ("Self-Rejecting Safety Filter", "Actively prune interventions that increase long-term cost or debt overhang."),
            ("Explainable Consent Layer", "SHAP feature attributions + Claude LLM plain-English customer explanations.")
        ], ACCENT_TEAL)
    ]

    for i, (heading, items, color) in enumerate(cards_data):
        left = Inches(0.8 + i * 5.9)
        add_card(slide3, left, Inches(1.85), Inches(5.7), Inches(4.9), CARD_BG, color)
        
        box = slide3.shapes.add_textbox(left + Inches(0.25), Inches(2.05), Inches(5.2), Inches(4.5))
        tf = box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = heading
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(14)
        
        for title, desc in items:
            p_t = tf.add_paragraph()
            p_t.text = ("✖  " if color == ACCENT_CORAL else "★  ") + title
            p_t.font.size = Pt(11)
            p_t.font.bold = True
            p_t.font.color.rgb = TEXT_WHITE
            
            p_d = tf.add_paragraph()
            p_d.text = "    " + desc
            p_d.font.size = Pt(10)
            p_d.font.color.rgb = TEXT_MUTED
            p_d.space_after = Pt(8)

    add_speaker_notes(slide3, """It's vital to understand what FinShield is NOT. We are not a credit scoring engine, we are not a loan upsell tool, and we are not a generic LLM chatbot.
FinShield operates like a clinical diagnostician. When a patient is sick, a doctor doesn't just assign a sickness score and walk away; they diagnose the specific ailment, measure the vital limits, and prescribe the gentlest, most effective treatment. That is 'Diagnosis Before Treatment'.""")

    # =========================================================================
    # SLIDE 4: THE END-TO-END FINSHIELD FLOW
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide4)
    add_header(slide4, "System Architecture", "The FinShield End-to-End Decision Flow", "A closed-loop pipeline combining machine learning, deterministic math, safety filters, and human oversight.")

    # 6 Step Cards arranged horizontally / vertically
    steps = [
        ("1. Early Distress Trigger", "Detects cash-flow strains, overdrafts, or expense spikes before payment default.", ACCENT_BLUE),
        ("2. Root-Cause Diagnosis", "XGBoost classifies distress: Income Shock, Debt Overload, Timing Mismatch, etc.", ACCENT_CYAN),
        ("3. Repayment Capacity", "Deterministic formula computes real monthly affordability floor (e.g. ₹15,000).", ACCENT_TEAL),
        ("4. Safety & Eligibility Filter", "Hard constraints reject dangerous options (e.g. rejects consolidation if debt-overloaded).", ACCENT_AMBER),
        ("5. Explainable Consent", "SHAP + Claude LLM present transparent options to customer with trade-offs.", ACCENT_CYAN),
        ("6. Safe Dual-Execution", "Tier A auto-executed post-consent; Tier B/C routed to Bank Hardship Officers.", ACCENT_TEAL)
    ]

    for i, (title, desc, col) in enumerate(steps):
        row = i // 3
        col_idx = i % 3
        left = Inches(0.8 + col_idx * 3.9)
        top = Inches(1.85 + row * 2.5)
        
        add_card(slide4, left, top, Inches(3.7), Inches(2.25), CARD_BG, CARD_BORDER)
        
        # Step Number Badge
        box = slide4.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), Inches(3.3), Inches(2.0))
        tf = box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = col
        p.space_after = Pt(6)
        
        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = TEXT_WHITE
        p_d.space_after = Pt(6)
        
        # Micro tag
        p_t = tf.add_paragraph()
        p_t.text = "→ Stage " + str(i+1) + " Output Validated"
        p_t.font.size = Pt(8.5)
        p_t.font.color.rgb = TEXT_DIM

    add_speaker_notes(slide4, """Here is the complete FinShield architecture in action.
From the moment high-risk signals appear, the system diagnoses the root cause, calculates sustainable repayment capacity, removes unsafe options through hard constraints, translates the plan via SHAP and Claude into plain English, secures customer consent, and executes safe actions while routing contractual changes to human officers.""")

    # =========================================================================
    # SLIDE 5: CORE AI & DECISION ENGINE STACK
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide5)
    add_header(slide5, "AI & Financial Engineering", "Deep Dive: The 4 Core Intelligence Modules", "AI is applied precisely where it adds distinct value — backed by deterministic financial guardrails.")

    modules = [
        ("A. XGBoost Distress Classifier", "Multi-Class ML Diagnosis", [
            "Predicts the underlying cause across 5 archetypes: Income Shock, Debt Overload, Cash-Flow Mismatch, Expense Shock, Structural Distress.",
            "Trained on multi-dimensional time-series features (cash-flow velocity, variance in discretionary vs essential spending, DTI trends)."
        ], ACCENT_BLUE),
        ("B. SHAP Explainability Engine", "Transparent Feature Attribution", [
            "Computes local Shapley values for every customer decision.",
            "Pinpoints exact drivers: e.g., +42% essential expense spike, -65% liquid savings drawdown, 3 overdraft occurrences.",
            "Zero black-box decisions — provides audit-ready regulatory transparency."
        ], ACCENT_CYAN),
        ("C. Repayment Capacity Engine", "Deterministic Affordability Math", [
            "Completely rule-based and transparent — NOT an opaque ML output.",
            "Formula: Monthly Net Income - Essential Living Floor - Buffer - Fixed Non-Bank Commitments = Max Sustainable EMI.",
            "Example: Current obligations ₹25k/mo → Sustainable floor ₹15k/mo."
        ], ACCENT_TEAL),
        ("D. Claude LLM Communication Layer", "Contextual Natural Language", [
            "Translates complex SHAP vectors and financial options into empathetic, jargon-free customer guidance.",
            "Generates comprehensive case dossiers for human hardship officers.",
            "Strict constraint: LLM has ZERO decision-making or rule-setting authority."
        ], RGBColor(168, 85, 247)) # Purple
    ]

    for i, (title, sub, bullets, col) in enumerate(modules):
        col_idx = i % 2
        row_idx = i // 2
        left = Inches(0.8 + col_idx * 5.9)
        top = Inches(1.85 + row_idx * 2.5)
        
        add_card(slide5, left, top, Inches(5.7), Inches(2.35), CARD_BG, CARD_BORDER)
        
        box = slide5.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), Inches(5.3), Inches(2.1))
        tf = box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = col
        
        p_sub = tf.add_paragraph()
        p_sub.text = sub.upper()
        p_sub.font.size = Pt(8.5)
        p_sub.font.bold = True
        p_sub.font.color.rgb = TEXT_DIM
        p_sub.space_after = Pt(4)
        
        for b in bullets:
            p_b = tf.add_paragraph()
            p_b.text = "• " + b
            p_b.font.size = Pt(9.5)
            p_b.font.color.rgb = TEXT_WHITE
            p_b.space_after = Pt(3)

    add_speaker_notes(slide5, """Let's look under the hood. We use AI where it excels, and deterministic logic where safety is paramount.
XGBoost categorizes the distress into one of 5 archetypes.
SHAP extracts the exact feature attributions explaining the diagnosis.
Crucially, the Repayment Capacity Engine is 100% deterministic math, not ML hallucinations.
Finally, Claude acts strictly as an empathetic translation layer — explaining the math in plain language, without having power to make financial approvals.""")

    # =========================================================================
    # SLIDE 6: THE 7-LEVEL INTERVENTION HIERARCHY
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide6)
    add_header(slide6, "Action Framework", "The 7-Level Intervention Ladder: Ordered by Intrusiveness", "Interventions are ranked strictly by reversibility and long-term cost, not simply customer risk score.")

    # 7 Levels arranged as vertical ladder
    ladder = [
        ("L0", "MONITOR", "Passive telemetry tracking when metrics are within safe variance limits. No customer contact.", ACCENT_CYAN, "Zero Friction"),
        ("L1", "INFORM", "Actionable spending insights, obligation calendar syncing, budgeting nudges.", ACCENT_CYAN, "Informational"),
        ("L2", "TIMING FIX", "Align repayment dates to actual payroll deposits, autopay re-timing, waiver of late fees.", ACCENT_TEAL, "High Reversibility"),
        ("L3", "REPAYMENT PLAN", "Structured 3-6 month catch-up schedule with stabilized monthly instalments.", ACCENT_TEAL, "Moderate Structure"),
        ("L4", "TEMPORARY RELIEF", "Short-term payment moratorium (1-3 months) or interest freeze during acute shock.", ACCENT_AMBER, "Relief Concession"),
        ("L5", "FORMAL SUPPORT", "Tenor extension, structured loan restructuring, contractual interest rate reduction.", ACCENT_AMBER, "Contractual Change"),
        ("L6", "HUMAN ESCALATION", "Specialized Hardship Officer assignment for complex, vulnerable, or disputed cases.", ACCENT_CORAL, "Human Judgment")
    ]

    for i, (lvl, name, desc, col, badge) in enumerate(ladder):
        top = Inches(1.8 + i * 0.72)
        
        # Row card
        add_card(slide6, Inches(0.8), top, Inches(11.73), Inches(0.65), CARD_BG, CARD_BORDER)
        
        # Pill badge
        badge_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95), top + Inches(0.12), Inches(0.8), Inches(0.4))
        badge_box.fill.solid()
        badge_box.fill.fore_color.rgb = col
        badge_box.line.fill.background()
        tf_b = badge_box.text_frame
        p_b = tf_b.paragraphs[0]
        p_b.text = lvl
        p_b.font.size = Pt(11)
        p_b.font.bold = True
        p_b.font.color.rgb = NAVY_DARK
        p_b.alignment = PP_ALIGN.CENTER
        
        # Name and description
        text_box = slide6.shapes.add_textbox(Inches(1.9), top + Inches(0.08), Inches(8.3), Inches(0.5))
        tf_t = text_box.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        p_t.text = name + " — "
        p_t.font.size = Pt(10.5)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_WHITE
        p_t.text += desc
        p_t.font.bold = False
        p_t.font.color.rgb = TEXT_MUTED
        
        # Trait Badge on Right
        trait_box = slide6.shapes.add_textbox(Inches(10.3), top + Inches(0.08), Inches(2.1), Inches(0.5))
        tf_tr = trait_box.text_frame
        p_tr = tf_tr.paragraphs[0]
        p_tr.text = badge.upper()
        p_tr.font.size = Pt(9)
        p_tr.font.bold = True
        p_tr.font.color.rgb = col
        p_tr.alignment = PP_ALIGN.RIGHT

    add_speaker_notes(slide6, """Most banking systems jump straight from doing nothing to aggressive collections or heavy loan restructuring.
FinShield introduces a 7-Level Intervention Ladder. It is strictly ordered by intrusiveness and reversibility.
Level 2 Timing Fixes—like simply adjusting the payment date to match when the customer gets paid—resolve over 30% of cash-flow friction without adding debt or requiring contractual restructuring!""")

    # =========================================================================
    # SLIDE 7: RESPONSIBLE AI & AUTOMATION BOUNDARY
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide7)
    add_header(slide7, "Governance & Ethics", "Responsible AI: Strict Automation Boundaries", "Responsible AI is not a disclaimer in our footer — it is hardcoded into our execution architecture.")

    # 3 Column Cards: Tier A, Tier B, Tier C
    tiers = [
        ("TIER A: AUTOMATE POST-CONSENT", ACCENT_TEAL, [
            "✔ Root-cause distress diagnosis",
            "✔ Sustainable capacity calculation",
            "✔ Reversible payment date alignment",
            "✔ Late-fee waivers within policy limits",
            "✔ Budgeting insights & SMS reminders",
            "✔ Continuous recovery telemetry"
        ], "Low-risk, 100% reversible actions with explicit customer consent."),
        
        ("TIER B: AI PREPARES, HUMAN APPROVES", ACCENT_AMBER, [
            "⚡ Temporary EMI moratorium (1-3 mos)",
            "⚡ Grace period activations",
            "⚡ Tenor extension proposals",
            "⚡ Interest rate concession packaging",
            "⚡ Comprehensive hardship case dossier",
            "⚡ Simulated outcome impact analysis"
        ], "Moderate financial changes; AI structures the proposal for loan officer sign-off."),
        
        ("TIER C: HUMAN ONLY (AI STRICTLY FORBIDDEN)", ACCENT_CORAL, [
            "✖ Denying banking services or credit",
            "✖ Involuntary loan restructuring",
            "✖ Adverse credit bureau reporting",
            "✖ Legal enforcement or repossession",
            "✖ Forcing customers into any intervention",
            "✖ Unilateral contractual alterations"
        ], "High-impact legal & contractual decisions remain 100% human-controlled.")
    ]

    for i, (title, col, items, footer) in enumerate(tiers):
        left = Inches(0.8 + i * 3.95)
        add_card(slide7, left, Inches(1.85), Inches(3.8), Inches(4.9), CARD_BG, col)
        
        box = slide7.shapes.add_textbox(left + Inches(0.2), Inches(2.0), Inches(3.4), Inches(4.6))
        tf = box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = col
        p.space_after = Pt(12)
        
        for item in items:
            p_i = tf.add_paragraph()
            p_i.text = item
            p_i.font.size = Pt(9.5)
            p_i.font.color.rgb = TEXT_WHITE
            p_i.space_after = Pt(5)
            
        p_foot = tf.add_paragraph()
        p_foot.text = "\n" + footer
        p_foot.font.size = Pt(9)
        p_foot.font.italic = True
        p_foot.font.color.rgb = TEXT_MUTED

    add_speaker_notes(slide7, """This is one of our most important slides. In banking, AI should never have unilateral authority to alter contractual terms, hurt a credit score, or deny services.
We established 3 strict operational tiers:
Tier A allows automated execution of low-risk, reversible actions only AFTER customer consent.
Tier B allows AI to assemble hardship packages, but a human officer must approve.
Tier C strictly prohibits AI from taking punitive, credit-reporting, or adverse legal actions.""")

    # =========================================================================
    # SLIDE 8: THE "WOW" MOMENT — SAFETY LAYER IN ACTION
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide8)
    add_header(slide8, "The Key Innovation", "The 'WOW' Moment: Rejecting Unsafe Recommendations", "FinShield doesn't just generate recommendations. It actively detects and rejects its own unsafe options.")

    # Left: The Rejection Box
    add_card(slide8, Inches(0.8), Inches(1.85), Inches(5.7), Inches(4.9), CARD_BG, ACCENT_CORAL)
    t_rej = slide8.shapes.add_textbox(Inches(1.05), Inches(2.05), Inches(5.2), Inches(4.5))
    tf_rej = t_rej.text_frame
    tf_rej.word_wrap = True
    
    p = tf_rej.paragraphs[0]
    p.text = "PROPOSED: DEBT CONSOLIDATION LOAN"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(8)
    
    p_badge = tf_rej.add_paragraph()
    p_badge.text = "❌ REJECTED BY FINSHIELD SAFETY ENGINE"
    p_badge.font.size = Pt(13)
    p_badge.font.bold = True
    p_badge.font.color.rgb = ACCENT_CORAL
    p_badge.space_after = Pt(12)
    
    eval_checks = [
        ("Customer Profile", "Debt Overload Archetype (DTI = 72%, Savings = 0.4 months)"),
        ("Hard Constraint Violated", "Rule SC-402: Cannot add/extend debt for borrowers with DTI > 65%"),
        ("Safety Audit Finding", "Consolidation would increase total lifetime interest by ₹84,000 and extend repayment duration without solving underlying cash deficit."),
        ("System Action", "Option permanently pruned from candidate set & logged to compliance audit trail.")
    ]
    for k, v in eval_checks:
        p_k = tf_rej.add_paragraph()
        p_k.text = "• " + k + ": "
        p_k.font.size = Pt(10)
        p_k.font.bold = True
        p_k.font.color.rgb = TEXT_WHITE
        p_k.text += v
        p_k.font.bold = False
        p_k.font.color.rgb = TEXT_MUTED
        p_k.space_after = Pt(6)

    # Right: Safe Alternative Generated
    add_card(slide8, Inches(6.8), Inches(1.85), Inches(5.7), Inches(4.9), CARD_BG, ACCENT_TEAL)
    t_safe = slide8.shapes.add_textbox(Inches(7.05), Inches(2.05), Inches(5.2), Inches(4.5))
    tf_safe = t_safe.text_frame
    tf_safe.word_wrap = True
    
    p = tf_safe.paragraphs[0]
    p.text = "SAFETY-FILTERED & RANKED ALTERNATIVE"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(8)
    
    p_badge = tf_safe.add_paragraph()
    p_badge.text = "✔ APPROVED: 2-STEP RECOVERY BUNDLE"
    p_badge.font.size = Pt(13)
    p_badge.font.bold = True
    p_badge.font.color.rgb = ACCENT_TEAL
    p_badge.space_after = Pt(12)
    
    safe_bundle = [
        ("Step 1: Payment-Date Sync (Level 2)", "Shift EMI due date from 1st to 7th (post-salary credit). Eliminates repeated ₹500 overdraft fees immediately."),
        ("Step 2: Temporary EMI Reduction (Level 4)", "Reduce monthly EMI from ₹25,000 to sustainable capacity floor of ₹15,000 for 3 months."),
        ("Net Outcome for Customer", "Instant cash-flow breathing room without taking on additional debt."),
        ("Net Outcome for Bank", "Cures arrears, prevents default write-off, maintains customer goodwill.")
    ]
    for k, v in safe_bundle:
        p_k = tf_safe.add_paragraph()
        p_k.text = "★ " + k + ": "
        p_k.font.size = Pt(10)
        p_k.font.bold = True
        p_k.font.color.rgb = ACCENT_CYAN
        p_k.text += v
        p_k.font.bold = False
        p_k.font.color.rgb = TEXT_WHITE
        p_k.space_after = Pt(6)

    add_speaker_notes(slide8, """This is our standout differentiator. Traditional banking systems love to recommend debt consolidation because it generates fee income.
When FinShield evaluates a debt-overloaded customer, the generation engine might produce a consolidation candidate. BUT our safety layer detects excessive debt burden and rejects it!
It audits WHY it was rejected and serves a safer, debt-free alternative instead.""")

    # =========================================================================
    # SLIDE 9: CUSTOMER EXPERIENCE JOURNEY (PRIYA'S STORY)
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide9)
    add_header(slide9, "User Experience", "Customer Journey: How Priya (Age 34) Experiences FinShield", "An empathetic, transparent mobile flow that restores customer dignity and control.")

    # 4 Mobile Mockup Cards
    journey_steps = [
        ("1. Early Proactive Alert", "Empathetic Outreach", [
            "“Hi Priya, we noticed an unexpected spike in medical expenses this month. Let’s keep your finances stress-free before your next EMI.”",
            "No punitive threats or threatening collection notices."
        ], ACCENT_BLUE),
        ("2. Transparent Diagnosis", "Clear Affordability Math", [
            "Expense Shock: 85% confidence",
            "Current Obligations: ₹25,000/mo",
            "Calculated Sustainable Floor: ₹15,000/mo",
            "Displays visual breakdown of essential vs debt obligations."
        ], ACCENT_CYAN),
        ("3. Choice & Comparison", "Tailored Options", [
            "Option A (Recommended): Due date shift + 3-mo ₹10k EMI reduction.",
            "Option B: 60-day principal pause.",
            "Option C: Connect with Hardship Specialist.",
            "Clear disclosure of trade-offs."
        ], ACCENT_TEAL),
        ("4. Consent & Recovery", "Real-Time Telemetry", [
            "Priya selects Option A with 1-click consent.",
            "Tier A date change executes instantly.",
            "Tier B temporary relief approved by officer.",
            "Dynamic tracker monitors financial recovery."
        ], ACCENT_AMBER)
    ]

    for i, (title, sub, bullets, col) in enumerate(journey_steps):
        left = Inches(0.8 + i * 2.95)
        add_card(slide9, left, Inches(1.85), Inches(2.8), Inches(4.9), CARD_BG, CARD_BORDER)
        
        box = slide9.shapes.add_textbox(left + Inches(0.15), Inches(2.0), Inches(2.5), Inches(4.6))
        tf = box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = col
        
        p_s = tf.add_paragraph()
        p_s.text = sub.upper()
        p_s.font.size = Pt(8)
        p_s.font.bold = True
        p_s.font.color.rgb = TEXT_DIM
        p_s.space_after = Pt(8)
        
        for b in bullets:
            p_b = tf.add_paragraph()
            p_b.text = "• " + b
            p_b.font.size = Pt(9)
            p_b.font.color.rgb = TEXT_WHITE
            p_b.space_after = Pt(4)

    add_speaker_notes(slide9, """Let's trace how Priya experiences FinShield. Priya earns ₹60,000/month with ₹25,000 in loan EMIs. A sudden medical emergency caused her expenses to surge.
Instead of sending threatening texts after she misses a payment, FinShield proactively notices the expense shock.
It shows her calculated sustainable capacity of ₹15,000, offers transparent options, obtains her consent, and monitors her recovery back to health.""")

    # =========================================================================
    # SLIDE 10: BANK OPERATIONS CONSOLE & WHAT-IF SIMULATOR
    # =========================================================================
    slide10 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide10)
    add_header(slide10, "Operations & Simulation", "Bank Operations Console & What-If Stress Simulator", "Empowering credit officers with portfolio-level telemetry and deterministic what-if projections.")

    # Left: Bank Ops Console
    add_card(slide10, Inches(0.8), Inches(1.85), Inches(5.7), Inches(4.9), CARD_BG, CARD_BORDER)
    t_ops = slide10.shapes.add_textbox(Inches(1.05), Inches(2.05), Inches(5.2), Inches(4.5))
    tf_ops = t_ops.text_frame
    tf_ops.word_wrap = True
    
    p = tf_ops.paragraphs[0]
    p.text = "PORTFOLIO OPERATIONS CONSOLE"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.space_after = Pt(10)
    
    ops_features = [
        ("Real-Time Triage Queue", "Categorizes distressed book by root-cause archetype and urgency."),
        ("Automated vs Officer Queue", "Tier A self-resolves; Tier B/C packages served with 1-click officer review."),
        ("SHAP Attribution Inspector", "Inspect individual feature weights and historical banking signals."),
        ("Regulatory Compliance Trail", "Full event log recording why interventions were selected or rejected.")
    ]
    for k, v in ops_features:
        p_k = tf_ops.add_paragraph()
        p_k.text = "✔ " + k + ": "
        p_k.font.size = Pt(10)
        p_k.font.bold = True
        p_k.font.color.rgb = TEXT_WHITE
        p_k.text += v
        p_k.font.bold = False
        p_k.font.color.rgb = TEXT_MUTED
        p_k.space_after = Pt(6)

    # Right: What-If Simulator
    add_card(slide10, Inches(6.8), Inches(1.85), Inches(5.7), Inches(4.9), CARD_BG, CARD_BORDER)
    t_sim = slide10.shapes.add_textbox(Inches(7.05), Inches(2.05), Inches(5.2), Inches(4.5))
    tf_sim = t_sim.text_frame
    tf_sim.word_wrap = True
    
    p = tf_sim.paragraphs[0]
    p.text = "WHAT-IF INTERVENTION SIMULATOR"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    p.space_after = Pt(10)
    
    sim_data = [
        ("BEFORE INTERVENTION (CRITICAL STRESS)", [
            "Monthly Income: ₹60,000 | Current EMI: ₹25,000",
            "Sustainable Capacity: ₹15,000 | Deficit: -₹10,000/mo",
            "Financial Stress Index: 88% (Severe Risk of Default)"
        ]),
        ("SIMULATED INTERVENTION: 3-MONTH EMI REDUCTION + DUE DATE SHIFT", [
            "Simulated Monthly EMI: ₹14,500 (Affordable)",
            "Projected Cash Flow Surplus: +₹1,500/mo buffer",
            "Projected Stress Index: 34% (Stable Recovery Zone)"
        ])
    ]
    for section_title, lines in sim_data:
        p_s = tf_sim.add_paragraph()
        p_s.text = section_title
        p_s.font.size = Pt(9.5)
        p_s.font.bold = True
        p_s.font.color.rgb = ACCENT_AMBER if "BEFORE" in section_title else ACCENT_TEAL
        for l in lines:
            p_l = tf_sim.add_paragraph()
            p_l.text = "  • " + l
            p_l.font.size = Pt(9)
            p_l.font.color.rgb = TEXT_WHITE
        p_l.space_after = Pt(6)
        
    p_disc = tf_sim.add_paragraph()
    p_disc.text = "⚠ TRANSPARENCY DISCLAIMER: SIMULATED ESTIMATE — NOT A GUARANTEED OUTCOME. Projections are rule-based feasibility estimates, not causal guarantees."
    p_disc.font.size = Pt(7.5)
    p_disc.font.bold = True
    p_disc.font.color.rgb = ACCENT_AMBER

    add_speaker_notes(slide10, """On the bank side, loan officers get an operational dashboard. It separates automated low-risk actions from high-impact cases requiring review.
The What-If Simulator allows officers and customers to test simulated interventions. Notice our transparent disclaimer: we explicitly label these as simulated estimates rather than claiming magical causal predictions.""")

    # =========================================================================
    # SLIDE 11: TECHNICAL STACK & PROTOTYPE METHODOLOGY
    # =========================================================================
    slide11 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide11)
    add_header(slide11, "Implementation & Engineering", "Technical Architecture & Prototype Methodology", "A high-performance, explainable fintech stack built for rapid deployment and strict compliance.")

    # 4 Stack Components
    stack = [
        ("Frontend & UX", "React 19, TailwindCSS, Recharts, Lucide", "Responsive portal with mobile-first customer journey, operations dashboard, interactive SHAP visualizations, and What-If simulator.", ACCENT_BLUE),
        ("Backend & API", "FastAPI, Python 3.13, Pydantic v2", "Asynchronous REST architecture, sub-50ms rule evaluation, deterministic capacity math engine, and idempotent action executors.", ACCENT_CYAN),
        ("Machine Learning & XAI", "XGBoost, SHAP TreeExplainer, Scikit-Learn", "Multi-class distress classifier trained on multi-variate financial time-series. Local SHAP values computed for each diagnosis.", ACCENT_TEAL),
        ("LLM & Data Infrastructure", "Claude 3.5 Sonnet API, PostgreSQL", "Generates natural language customer coaching and officer hardship dossiers. PostgreSQL event store for immutable compliance audit logs.", RGBColor(168, 85, 247))
    ]

    for i, (title, tech, desc, col) in enumerate(stack):
        col_idx = i % 2
        row_idx = i // 2
        left = Inches(0.8 + col_idx * 5.9)
        top = Inches(1.85 + row_idx * 2.1)
        
        add_card(slide11, left, top, Inches(5.7), Inches(1.95), CARD_BG, CARD_BORDER)
        
        box = slide11.shapes.add_textbox(left + Inches(0.2), top + Inches(0.12), Inches(5.3), Inches(1.75))
        tf = box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = col
        
        p_t = tf.add_paragraph()
        p_t.text = tech
        p_t.font.size = Pt(9.5)
        p_t.font.bold = True
        p_t.font.color.rgb = TEXT_WHITE
        p_t.space_after = Pt(3)
        
        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(9)
        p_d.font.color.rgb = TEXT_MUTED

    # Bottom Banner: Synthetic Data Transparency Statement
    add_card(slide11, Inches(0.8), Inches(6.15), Inches(11.73), Inches(0.7), CARD_BG, ACCENT_AMBER)
    b_box = slide11.shapes.add_textbox(Inches(0.95), Inches(6.18), Inches(11.4), Inches(0.6))
    tf_b = b_box.text_frame
    tf_b.word_wrap = True
    p_b = tf_b.paragraphs[0]
    p_b.text = "★ DATASET TRANSPARENCY: "
    p_b.font.size = Pt(9)
    p_b.font.bold = True
    p_b.font.color.rgb = ACCENT_AMBER
    p_b.text += "Our hackathon prototype is validated on a rigorously engineered synthetic time-series dataset modeling 5 distress archetypes with realistic income, expense, and payment shocks. We do not claim validation on proprietary live bank production data."
    p_b.font.bold = False
    p_b.font.color.rgb = TEXT_MUTED

    add_speaker_notes(slide11, """Here is our technical stack: FastAPI backend, XGBoost and SHAP for ML, Claude API for natural language translation, and React with Tailwind and Recharts on the frontend.
In the spirit of scientific integrity, we clearly state that our prototype is evaluated on a synthetic financial time-series dataset simulating the 5 distress archetypes, rather than claiming untested live bank validation.""")

    # =========================================================================
    # SLIDE 12: QUALITATIVE IMPACT & CLOSING VISION
    # =========================================================================
    slide12 = prs.slides.add_slide(blank_layout)
    set_slide_background(slide12)
    add_header(slide12, "Strategic Impact & Closing", "The Paradigm Shift: From Detection to Safe Resolution", "Don't wait for default to act. Transform credit operations from punitive recovery to proactive care.")

    # Left: Impact for Customers & Banks
    add_card(slide12, Inches(0.8), Inches(1.85), Inches(5.7), Inches(3.6), CARD_BG, CARD_BORDER)
    t_imp = slide12.shapes.add_textbox(Inches(1.0), Inches(1.95), Inches(5.3), Inches(3.4))
    tf_imp = t_imp.text_frame
    tf_imp.word_wrap = True
    
    p = tf_imp.paragraphs[0]
    p.text = "MEASURABLE QUALITATIVE VALUE"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.space_after = Pt(8)
    
    impacts = [
        ("For Borrowers", "Proactive assistance, realistic affordability, preserved dignity, and avoidance of damaging debt spirals."),
        ("For Financial Institutions", "Higher cure rates on distressed books, 60% reduction in manual triage time, and automated compliance audit trails."),
        ("For Regulators", "Guaranteed adherence to fair lending laws, explainable AI decisions, and zero opaque black-box debt traps.")
    ]
    for k, v in impacts:
        p_k = tf_imp.add_paragraph()
        p_k.text = "✔ " + k + ": "
        p_k.font.size = Pt(9.5)
        p_k.font.bold = True
        p_k.font.color.rgb = TEXT_WHITE
        p_k.text += v
        p_k.font.bold = False
        p_k.font.color.rgb = TEXT_MUTED
        p_k.space_after = Pt(4)

    # Right: The Paradigm Shift
    add_card(slide12, Inches(6.8), Inches(1.85), Inches(5.7), Inches(3.6), CARD_BG, CARD_BORDER)
    t_para = slide12.shapes.add_textbox(Inches(7.0), Inches(1.95), Inches(5.3), Inches(3.4))
    tf_para = t_para.text_frame
    tf_para.word_wrap = True
    
    p = tf_para.paragraphs[0]
    p.text = "THE OPERATIONAL PARADIGM SHIFT"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT_TEAL
    p.space_after = Pt(8)
    
    p_old = tf_para.add_paragraph()
    p_old.text = "TRADITIONAL MODEL (PUNITIVE & LATE):\n"
    p_old.font.size = Pt(9.5)
    p_old.font.bold = True
    p_old.font.color.rgb = ACCENT_CORAL
    p_old.text += "Detect (Late)  ➔  Collect (Aggressive)  ➔  Escalate (Default & Loss)"
    p_old.font.bold = False
    p_old.font.color.rgb = TEXT_MUTED
    p_old.space_after = Pt(12)
    
    p_new = tf_para.add_paragraph()
    p_new.text = "FINSHIELD MODEL (PREVENTATIVE & SAFE):\n"
    p_new.font.size = Pt(9.5)
    p_new.font.bold = True
    p_new.font.color.rgb = ACCENT_TEAL
    p_new.text += "Diagnose (Early)  ➔  Intervene (Safely)  ➔  Recover (Sustainable)"
    p_new.font.bold = False
    p_new.font.color.rgb = TEXT_WHITE

    # Bottom Full Width Callout Banner
    add_card(slide12, Inches(0.8), Inches(5.65), Inches(11.73), Inches(1.2), CARD_BG, ACCENT_BLUE)
    t_call = slide12.shapes.add_textbox(Inches(1.0), Inches(5.7), Inches(11.3), Inches(1.1))
    tf_call = t_call.text_frame
    tf_call.word_wrap = True
    
    p_c = tf_call.paragraphs[0]
    p_c.text = "“AI doesn't replace human judgment. It makes responsible intervention faster, safer, and more personalized.”"
    p_c.font.size = Pt(14)
    p_c.font.bold = True
    p_c.font.color.rgb = TEXT_WHITE
    p_c.alignment = PP_ALIGN.CENTER
    p_c.space_after = Pt(4)
    
    p_c2 = tf_call.add_paragraph()
    p_c2.text = "FINSHIELD  |  DON'T WAIT FOR DEFAULT TO ACT"
    p_c2.font.size = Pt(10)
    p_c2.font.bold = True
    p_c2.font.color.rgb = ACCENT_CYAN
    p_c2.alignment = PP_ALIGN.CENTER

    add_speaker_notes(slide12, """To conclude: Banks can no longer afford to wait for default before acting.
FinShield shifts banking from 'Detect, Collect, Escalate' to 'Diagnose, Intervene, Recover'.
By pairing machine learning diagnosis with deterministic safety checks and human oversight, FinShield protects both the bank's balance sheet and the customer's financial dignity.
Thank you, and we look forward to your questions!""")

    # Save presentation
    prs.save(filename)
    print(f"Presentation saved successfully to {filename}")

if __name__ == "__main__":
    out_file = sys.argv[1] if len(sys.argv) > 1 else "FinShield_Hackathon_Pitch.pptx"
    create_deck(out_file)
